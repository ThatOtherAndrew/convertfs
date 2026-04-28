from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from convertfs.converters.markitdown import MarkItDownDocuments


def _write_minimal_pptx(path: Path, text: str) -> None:
	presentation = Presentation()
	slide = presentation.slides.add_slide(presentation.slide_layouts[5])
	textbox = slide.shapes.add_textbox(0, 0, presentation.slide_width, presentation.slide_height)
	textbox.text_frame.text = text
	presentation.save(path)


def test_markitdown_documents_converts_pptx_to_markdown(tmp_path: Path) -> None:
	source = tmp_path / 'hello.pptx'
	_write_minimal_pptx(source, 'Hello World')

	dest = tmp_path / 'out.md'
	MarkItDownDocuments().process(source, tmp_path / 'hello.md', dest)
	result = dest.read_text(encoding='utf-8')

	assert 'Hello World' in result
	assert result.strip()